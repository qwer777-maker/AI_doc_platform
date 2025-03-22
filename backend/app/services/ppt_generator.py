from typing import Dict, Any, Optional, List
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.text.text import TextFrame
from pptx.enum.shapes import MSO_SHAPE
import logging

from .ai_service_factory import AIServiceFactory

# 配置日志
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class PPTGenerator:
    def __init__(self, ai_service_type: str = "deepseek"):
        self.templates_dir = os.path.join(os.path.dirname(__file__), "../templates/ppt_templates")
        # 修正路径：生成文档目录和app是同级的
        self.output_dir = os.path.abspath(os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), "generated_docs"))
        logger.info(f"PPT生成器输出目录: {self.output_dir}")
        os.makedirs(self.output_dir, exist_ok=True)
        
        # 确保目录权限正确
        try:
            os.chmod(self.output_dir, 0o777)
        except Exception as e:
            logger.warning(f"无法修改目录权限: {e}")
        
        # 定义配色方案
        self.COLORS = {
            'primary': RGBColor(41, 128, 185),    # 主色调：蓝色
            'secondary': RGBColor(44, 62, 80),    # 次要色：深灰
            'accent': RGBColor(231, 76, 60),      # 强调色：红色
            'light': RGBColor(236, 240, 241),     # 浅色：近白
            'dark': RGBColor(52, 73, 94)          # 深色：深灰蓝
        }
        
        self.prs = None
        self.ai_service = AIServiceFactory.create_service(ai_service_type)

    def generate(self, topic: str, outline: List[Dict[str, Any]], template_id: Optional[str] = None) -> Optional[str]:
        """
        生成PPT文档
        
        Args:
            topic: 主题
            outline: 大纲内容
            template_id: 模板ID（现在不使用，保留参数以保持接口兼容）
        """
        try:
            logger.info(f"开始生成PPT: 主题='{topic}', 章节数={len(outline)}")
            if template_id:
                logger.info(f"使用模板: {template_id}")
            
            # 创建新的演示文稿
            self.prs = Presentation()
            self.prs.slide_width = Inches(16)
            self.prs.slide_height = Inches(9)
            logger.info(f"创建新的演示文稿: 宽度={self.prs.slide_width}, 高度={self.prs.slide_height}")
            
            # 添加标题幻灯片
            logger.info("添加标题幻灯片")
            self._add_title_slide(topic)
            
            # 添加目录幻灯片
            logger.info("添加目录幻灯片")
            self._add_toc_slide(outline)
            
            # 处理每个章节
            total_slides = 2  # 已添加标题和目录幻灯片
            for section_index, section in enumerate(outline):
                section_title = section["title"]
                logger.info(f"处理章节 {section_index+1}/{len(outline)}: '{section_title}'")
                
                # 添加章节标题幻灯片
                self._add_section_title_slide(section_title)
                total_slides += 1
                
                # 添加章节内容幻灯片
                slides = section.get("slides", [])
                logger.info(f"  章节 '{section_title}' 包含 {len(slides)} 张幻灯片")
                
                for slide_index, slide_content in enumerate(slides):
                    slide_title = slide_content.get("title", "未知标题")
                    slide_type = slide_content.get("type", "content")
                    logger.info(f"  添加幻灯片 {slide_index+1}/{len(slides)}: '{slide_title}' (类型: {slide_type})")
                    self._add_content_slide(slide_content, topic, section_title)
                    total_slides += 1
            
            # 添加结束幻灯片
            logger.info("添加结束幻灯片")
            self._add_ending_slide(topic)
            total_slides += 1
            
            # 保存文件
            output_path = os.path.join(self.output_dir, f"{topic}_presentation.pptx")
            self.prs.save(output_path)
            
            logger.info(f"PPT生成完成: 共 {total_slides} 张幻灯片, 保存至: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"生成PPT时出错: {str(e)}")
            return None
        finally:
            self.prs = None

    def _add_points(self, text_frame: TextFrame, points: List[Dict[str, Any]]) -> None:
        """添加要点和详细说明，自动调整字体大小和布局"""
        # 要点数量的阈值，超过此数量则减小字体大小
        points_threshold = 3
        details_threshold = 2
        
        # 根据要点数量调整字体大小
        main_point_size = Pt(28)
        detail_point_size = Pt(20)
        
        if len(points) > points_threshold:
            # 要点较多，减小字体
            main_point_size = Pt(24)
            detail_point_size = Pt(18)
            
        for i, point in enumerate(points):
            # 添加主要要点
            p = text_frame.add_paragraph()
            p.text = "▪ " + point.get("main", "")  # 添加项目符号
            p.level = 0
            p.font.name = '微软雅黑'
            p.font.size = main_point_size
            p.font.bold = True
            p.font.color.rgb = self.COLORS['primary']
            p.space_after = Pt(8) if len(points) > points_threshold else Pt(12)  # 根据内容量调整间距
            
            # 获取详细说明
            details = point.get("details", [])
            
            # 如果详细说明太多，减小字体大小和间距
            detail_font_size = detail_point_size
            if len(details) > details_threshold:
                detail_font_size = Pt(16)
            
            # 添加详细说明
            for detail in details:
                p = text_frame.add_paragraph()
                p.text = "• " + detail
                p.level = 1
                p.font.name = '微软雅黑'
                p.font.size = detail_font_size
                p.font.color.rgb = self.COLORS['secondary']
                p.space_before = Pt(4) if len(details) > details_threshold else Pt(6)  # 根据内容量调整间距
                p.space_after = Pt(4) if len(details) > details_threshold else Pt(6)   # 根据内容量调整间距

    def _add_content_slide(self, content: Dict[str, Any], topic: str, section_title: str) -> None:
        """添加内容幻灯片，根据内容量自动调整布局"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 添加标题和装饰线条
        title = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(14), Inches(1)
        )
        title.text_frame.text = content.get("title", "")
        p = title.text_frame.paragraphs[0]
        p.font.name = '微软雅黑'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['dark']
        p.alignment = PP_ALIGN.LEFT
        
        # 添加装饰线条（使用矩形作为线条）
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(1.3), Inches(14), Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.COLORS['primary']
        
        # 获取幻灯片内容
        slide_title = content.get("title", "")
        slide_type = content.get("type", "content")
        
        # 使用AI服务生成幻灯片内容
        slide_content = self.ai_service.generate_slide_content(topic, section_title, slide_title, slide_type)
        
        # 获取要点数量，判断是否需要双列布局
        points = slide_content.get("points", [])
        total_points = len(points)
        total_details = sum(len(point.get("details", [])) for point in points)
        
        # 自动判断布局方式
        use_two_column = total_points > 3 or total_details > 6
        
        if slide_type == "two_column" or use_two_column:
            # 如果指定双列或内容较多，使用两栏布局
            if slide_type == "two_column" and "left_points" in slide_content and "right_points" in slide_content:
                # 使用预定义的左右栏内容
                left_content = slide.shapes.add_textbox(
                    Inches(1), Inches(1.5), Inches(6.5), Inches(6)
                )
                right_content = slide.shapes.add_textbox(
                    Inches(8), Inches(1.5), Inches(6.5), Inches(6)
                )
                self._add_points(left_content.text_frame, slide_content.get("left_points", []))
                self._add_points(right_content.text_frame, slide_content.get("right_points", []))
            else:
                # 自动将内容分为左右两栏
                left_points = points[:len(points)//2 + len(points)%2]  # 左侧放置一半+余数
                right_points = points[len(points)//2 + len(points)%2:]  # 右侧放置剩余部分
                
                left_content = slide.shapes.add_textbox(
                    Inches(1), Inches(1.5), Inches(6.5), Inches(6)
                )
                right_content = slide.shapes.add_textbox(
                    Inches(8), Inches(1.5), Inches(6.5), Inches(6)
                )
                self._add_points(left_content.text_frame, left_points)
                self._add_points(right_content.text_frame, right_points)
        else:
            # 普通单列布局
            text_content = slide.shapes.add_textbox(
                Inches(1), Inches(1.5), 
                Inches(14 if slide_type == "content" else 8), 
                Inches(6.5)  # 增加高度
            )
            self._add_points(text_content.text_frame, points)
            
            if slide_type == "image_content":
                # 在这里我们只添加图片描述，实际项目中可以根据描述生成或选择图片
                image_desc = slide_content.get("image_description", "")
                if image_desc:
                    image_note = slide.shapes.add_textbox(
                        Inches(9.5), Inches(3.5), Inches(5.5), Inches(1)
                    )
                    image_note.text_frame.text = f"[图片: {image_desc}]"
                    p = image_note.text_frame.paragraphs[0]
                    p.font.italic = True
                    p.font.color.rgb = self.COLORS['secondary']
        
        # 添加注释
        if notes := content.get("notes"):
            slide.notes_slide.notes_text_frame.text = notes

    def _add_title_slide(self, topic: str) -> None:
        """添加标题幻灯片"""
        slide_layout = self.prs.slide_layouts[0]  # 使用标题布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = topic
        subtitle.text = "AI自动生成的演示文稿"
        
        # 设置标题样式
        for paragraph in title.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(44)
                run.font.bold = True
                run.font.color.rgb = self.COLORS['primary']
        
        # 设置副标题样式
        for paragraph in subtitle.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(24)
                run.font.italic = True
                run.font.color.rgb = self.COLORS['secondary']

    def _add_toc_slide(self, outline: List[Dict[str, Any]]) -> None:
        """添加目录幻灯片"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 使用空白布局
        
        # 添加标题
        title = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(14), Inches(1)
        )
        title.text_frame.text = "目录"
        title_p = title.text_frame.paragraphs[0]
        title_p.alignment = PP_ALIGN.CENTER
        title_p.font.size = Pt(40)
        title_p.font.bold = True
        title_p.font.color.rgb = self.COLORS['primary']
        
        # 添加装饰线条
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,  # 修改这里
            Inches(2), Inches(1.8), Inches(12), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.COLORS['primary']
        
        # 创建两列布局
        left_content = slide.shapes.add_textbox(
            Inches(2), Inches(2.2), Inches(5.5), Inches(5)
        )
        right_content = slide.shapes.add_textbox(
            Inches(8.5), Inches(2.2), Inches(5.5), Inches(5)
        )
        
        # 分配目录项到两列
        mid_point = len(outline) // 2 + len(outline) % 2
        
        # 添加左列目录项
        tf_left = left_content.text_frame
        for i, section in enumerate(outline[:mid_point], 1):
            p = tf_left.add_paragraph()
            p.text = f"{i}. {section['title']}"
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['secondary']
            p.space_after = Pt(20)  # 增加段落间距
        
        # 添加右列目录项
        tf_right = right_content.text_frame
        for i, section in enumerate(outline[mid_point:], mid_point + 1):
            p = tf_right.add_paragraph()
            p.text = f"{i}. {section['title']}"
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['secondary']
            p.space_after = Pt(20)  # 增加段落间距

    def _add_section_title_slide(self, section_title: str) -> None:
        """添加章节标题幻灯片"""
        slide_layout = self.prs.slide_layouts[2]  # 使用章节标题布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = section_title
        
        # 设置标题样式
        for paragraph in title.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(40)
                run.font.bold = True
                run.font.color.rgb = self.COLORS['primary']

    def _add_ending_slide(self, topic: str) -> None:
        """添加结束幻灯片"""
        slide_layout = self.prs.slide_layouts[1]  # 使用标题和内容布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "总结与问答"
        
        # 设置标题样式
        for paragraph in title.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(40)
                run.font.bold = True
                run.font.color.rgb = self.COLORS['primary']
        
        # 添加总结内容
        tf = content.text_frame
        
        summary_points = [
            f"我们探讨了{topic}的多个关键方面",
            f"理解这些概念对把握{topic}的本质至关重要",
            f"未来发展将带来更多机遇和挑战",
            f"感谢您的关注！有任何问题欢迎提问"
        ]
        
        for point in summary_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            
            # 设置段落样式
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(28)
                run.font.color.rgb = self.COLORS['secondary'] 